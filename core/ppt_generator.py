"""
Premium PowerPoint Generator Utility
Creates stunning, tech-first presentations with a futuristic dark aesthetic.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import os

# Color Palette (Standard Meridian Tokens)
COLOR_BG_DARK = RGBColor(10, 10, 12)
COLOR_TEXT_PRIMARY = RGBColor(248, 250, 252) # Slate 50
COLOR_TEXT_SECONDARY = RGBColor(148, 163, 184) # Slate 400
COLOR_ACCENT_BLUE = RGBColor(10, 132, 255) # Apple Blue
COLOR_ACCENT_PURPLE = RGBColor(139, 92, 246) # Gemini Purple

ASSETS_DIR = "/Volumes/BLACK_SHARK/MINOR PROJECT/static/assets"
BG_IMAGE_PATH = os.path.join(ASSETS_DIR, "meridian_bg.png")

class PPTGenerator:
    def __init__(self, title="Meridian Data Report"):
        self.prs = Presentation()
        # Set 16:9 slide size
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.title = title
        
    def _add_background(self, slide):
        """Adds a premium background image and brand overlay to a slide."""
        if os.path.exists(BG_IMAGE_PATH):
            slide.shapes.add_picture(BG_IMAGE_PATH, 0, 0, width=self.prs.slide_width, height=self.prs.slide_height)
        
        # Add Brand Text Corner (Top Right)
        brand_box = slide.shapes.add_textbox(Inches(10.5), Inches(0.2), Inches(2.5), Inches(0.5))
        tf = brand_box.text_frame
        p = tf.paragraphs[0]
        p.text = "✨ Meridian Data"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        p.alignment = PP_ALIGN.RIGHT

    def _add_glass_panel(self, slide, left, top, width, height):
        """Adds a semi-transparent 'glass' panel effect."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(24, 24, 27) # Darker card color
        fill.transparency = 0.4
        
        line = shape.line
        line.color.rgb = RGBColor(255, 255, 255)
        line.width = Pt(0.5)
        line.transparency = 0.8
        
        return shape

    def add_title_slide(self, subtitle="Intelligent Data Analysis"):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6]) # Blank layout
        self._add_background(slide)
        
        # Central Glow / Content Box
        panel_w, panel_h = Inches(8), Inches(3)
        left, top = (self.prs.slide_width - panel_w) / 2, (self.prs.slide_height - panel_h) / 2
        self._add_glass_panel(slide, left, top, panel_w, panel_h)
        
        # Title
        title_box = slide.shapes.add_textbox(left, top + Inches(0.5), panel_w, Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self.title
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # Accent Line
        line_w = Inches(3)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, (self.prs.slide_width - line_w) / 2, top + Inches(1.6), line_w, Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT_PURPLE
        line.line.visible = False
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(left, top + Inches(1.8), panel_w, Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = COLOR_TEXT_SECONDARY
        p.alignment = PP_ALIGN.CENTER

    def add_text_slide(self, title, content):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        
        # Title Ribbon
        self._add_glass_panel(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.6))
        title_box.text_frame.text = title
        p = title_box.text_frame.paragraphs[0]
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        
        # Content Panel
        self._add_glass_panel(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.4), Inches(4.8))
        body_box.text_frame.word_wrap = True
        
        # Split content into bullets if it looks like bullet points
        lines = content.split('\n')
        for i, line in enumerate(lines):
            if not line.strip(): continue
            p = body_box.text_frame.add_paragraph() if i > 0 else body_box.text_frame.paragraphs[0]
            p.text = line.strip().lstrip('- •')
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_TEXT_PRIMARY
            p.space_after = Pt(10)
            if line.strip().startswith(('-', '•')):
                p.level = 1

    def add_table_slide(self, title, columns, rows):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        
        # Title
        self._add_glass_panel(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.6))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(28)
        title_box.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_PRIMARY
        
        # Table Container
        self._add_glass_panel(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        
        rows_count = min(len(rows) + 1, 11)
        cols_count = len(columns)
        
        left, top = Inches(0.8), Inches(1.8)
        width, height = Inches(11.4), Inches(0.4 * rows_count)
        
        shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
        table = shape.table
        
        # Style Headers
        for i, col in enumerate(columns):
            cell = table.cell(0, i)
            cell.text = str(col).upper()
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 41, 59) # Slate 800
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_ACCENT_BLUE
            
        # Style Data
        for r_idx, row in enumerate(rows[:10]):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.color.rgb = COLOR_TEXT_PRIMARY
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate 900

    def add_chart_slide(self, title, chart_config):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_background(slide)
        
        # Title
        self._add_glass_panel(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(10), Inches(0.6))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(28)
        title_box.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT_PRIMARY
        
        # Chart data
        chart_data = CategoryChartData()
        chart_data.categories = chart_config.get("labels", [])
        for ds in chart_config.get("datasets", []):
            chart_data.add_series(ds.get("label", ""), ds.get("data", []))
            
        # Chart Panel
        self._add_glass_panel(slide, Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
        }
        ct = chart_type_map.get(chart_config.get("type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
        
        # Positioning
        x, y, cx, cy = Inches(1), Inches(1.8), Inches(11), Inches(4.8)
        chart_shape = slide.shapes.add_chart(ct, x, y, cx, cy, chart_data)
        chart = chart_shape.chart
        
        # Style Chart (Some basics, python-pptx chart styling is complex)
        chart.has_legend = True
        chart.legend.font.color.rgb = COLOR_TEXT_SECONDARY
        chart.chart_title.has_text_frame = False # We have slide titles
        
    def save(self):
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output
