#!/usr/bin/env python3
"""Fill the CS3905 mini-project PPT template with Meridian Data content.
Preserves template formatting (clones paragraph XML); keeps text above the footer band."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.enum.text import MSO_ANCHOR
from PIL import Image

SRC = "/Users/ayonaryan/Downloads/MP/Phase 3_Mini Project ( CS3905).pptx"
OUT = "/Users/ayonaryan/Downloads/MP/Meridian_Data_Phase3_Presentation.pptx"
DIAG = "/Volumes/BLACK_SHARK/MINOR_PROJECT/report_build/assets/diagrams"
SCR = "/Volumes/BLACK_SHARK/MINOR_PROJECT/report_build/assets/screens"

prs = Presentation(SRC)

def banner_shape(slide):
    for sh in slide.shapes:
        if sh.shape_type == 1 and sh.has_text_frame and sh.text_frame.text.strip():
            return sh

def body_shape(slide):
    cand = [sh for sh in slide.shapes if sh.shape_type == 17 and sh.has_text_frame]
    return max(cand, key=lambda s: (s.width or 0)*(s.height or 0)) if cand else None

def set_banner(slide, title):
    sh = banner_shape(slide); para = sh.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = title
        for r in para.runs[1:]: r._r.getparent().remove(r._r)

def _first_run_para(tf):
    for p in tf.paragraphs:
        if p.runs: return p._p
    return tf.paragraphs[0]._p

def _runs_for(item):
    if isinstance(item, tuple):
        return [(item[0], item[1])]
    if " — " in item:                      # bold the lead-in before em dash
        lead, rest = item.split(" — ", 1)
        return [(lead + " — ", True), (rest, False)]
    return [(item, False)]

def set_bullets(shape, items, size=None):
    tf = shape.text_frame
    donor = copy.deepcopy(_first_run_para(tf))
    rtmpl = donor.find(qn('a:r'))
    for r in donor.findall(qn('a:r')): donor.remove(r)
    body = tf._txBody
    for p in body.findall(qn('a:p')): body.remove(p)
    for item in items:
        p = copy.deepcopy(donor)
        for text, bold in _runs_for(item):
            r = copy.deepcopy(rtmpl)
            t = r.find(qn('a:t'))
            if t is None:
                t = r.makeelement(qn('a:t'), {}); r.append(t)
            t.text = text
            rPr = r.find(qn('a:rPr'))
            if rPr is None:
                rPr = r.makeelement(qn('a:rPr'), {}); r.insert(0, rPr)
            rPr.set('b', '1' if bold else '0')
            if size is not None:
                rPr.set('sz', str(int(size*100)))
            p.append(r)
        body.append(p)

def configure_body(shape, top=1.45, height=3.55, left=None, width=None):
    shape.top = Inches(top); shape.height = Inches(height)
    if left is not None: shape.left = Inches(left)
    if width is not None: shape.width = Inches(width)
    tf = shape.text_frame
    tf.word_wrap = True
    try: tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception: pass

def fit(path, bw, bh):
    w, h = Image.open(path).size
    sc = min(bw/(w/96.0), bh/(h/96.0))
    return Inches((w/96.0)*sc), Inches((h/96.0)*sc)

def add_image(slide, path, bl, bt, bw, bh):
    iw, ih = fit(path, bw, bh)
    return slide.shapes.add_picture(path, Inches(bl)+(Inches(bw)-iw)//2, Inches(bt)+(Inches(bh)-ih)//2, iw, ih)

def remove(shape):
    shape._element.getparent().remove(shape._element)

def duplicate(prs, idx):
    src = prs.slides[idx]
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes): remove(sh)
    for sh in src.shapes: new.shapes._spTree.append(copy.deepcopy(sh._element))
    return new

s = list(prs.slides)

# ---- Slide 1: Title ----
t = s[0]; tb = t.shapes[0].text_frame
tb.paragraphs[1].runs[0].text = "Meridian Data"
p2 = tb.paragraphs[2]; p2.runs[0].text = "An AI-Powered NL-to-SQL Database Explorer"
for r in p2.runs[1:]: r._r.getparent().remove(r._r)
set_bullets(t.shapes[4], ["Team Members", "1RVU23CSE093  —  Ayon Aryan",
    "1RVU23CSE180  —  Hardik Goel", "1RVU23CSE055  —  Aniket Kumar",
    "1RVU23CSE029  —  Aditya Kumar"])
set_bullets(t.shapes[5], ["Project Guide:", "Prof. Ashwini Mathur",
    "Assistant Professor", "School of Computer Science and Engineering"])
for sh in t.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().lower().startswith("internal"):
        sh.text_frame.paragraphs[0].runs[0].text = "Internal"

# ---- Slide 2: Contents ----
configure_body(body_shape(s[1]), top=1.25, height=3.9)
set_bullets(body_shape(s[1]), [
    "Title & Problem Definition", "Introduction & Scope", "Objectives",
    "Literature Survey", "Proposed Methodology", "Requirement Specification",
    "Design Approach & Architecture", "Innovations", "Implementation",
    "Evaluation, Testing & Performance", "Conclusion & Future Work", "References",
])

# ---- Slide 3: Problem Definition ----
configure_body(body_shape(s[2]))
set_bullets(body_shape(s[2]), [
    "Problem — querying databases requires SQL fluency, a barrier for students, analysts and domain experts.",
    "Existing systems — single proprietary engine, closed-source, and send data to the cloud unconditionally.",
    "They hallucinate SQL and lack safety, authorization and human oversight.",
    "Novelty — hybrid deterministic + LLM querying over eight engines, with layered safety, human review and cloud-to-local privacy.",
])

# ---- Slide 4: Introduction & Scope ----
configure_body(body_shape(s[3]))
set_bullets(body_shape(s[3]), [
    "Converts plain-English requests into validated, dialect-specific SQL using LLMs given the live schema.",
    "Scope — NL querying, deterministic introspection, AI analysis with auto-charts, dashboards, and CSV / PowerPoint export.",
    "Connects to SQLite, PostgreSQL, MySQL, MSSQL, Oracle, MongoDB, Cassandra and Redis via one adapter layer.",
    "Doubles as a productivity tool and a DBMS teaching aid — it shows and explains the SQL it generates.",
])

# ---- Slide 5: Objectives ----
configure_body(body_shape(s[4]))
set_bullets(body_shape(s[4]), [
    "Convert natural language into validated, dialect-aware SQL via a schema-aware prompt.",
    "Answer common introspection commands deterministically, without the LLM.",
    "Enforce two-layer safety with role-based access control.",
    "Keep a human in the loop for writes — dry-run, snapshot and rollback.",
    "Provide AI analysis, auto chart selection, dashboards and export.",
    "Support cloud (Groq) and local (Ollama) models with automatic failover.",
])

# ---- Slide 6: Literature Survey ----
configure_body(body_shape(s[5]))
set_bullets(body_shape(s[5]), [
    "Evolved from rule-based NLIDBs to neural text-to-SQL (Seq2SQL / WikiSQL, SQLNet) and the Spider benchmark.",
    "RAT-SQL and BRIDGE solved schema linking; PICARD added constrained decoding for valid SQL.",
    "LLM era — in-context learning, DIN-SQL decomposition, DAIL-SQL prompting; BIRD shows real databases stay hard (GPT-4 ~55%).",
    "Safety / RBAC and human-in-the-loop research inform the design.",
    "Meridian Data integrates these validated components rather than proposing a new parser.",
])

# ---- Slide 7: Proposed Methodology ----
configure_body(body_shape(s[6]))
set_bullets(body_shape(s[6]), [
    "Intent classification — introspection answered deterministically; data questions go to the LLM.",
    "Schema-aware generation — schema (columns, keys, indexes, samples) injected into a dialect prompt; Groq / Ollama failover.",
    "Validation & authorization — classify (READ / WRITE / SCHEMA / SYSTEM) → safety filter → role check.",
    "Execution or review — READ runs paginated; WRITE / SCHEMA go to human review with snapshot.",
    "Analysis & presentation — AI summary, auto-selected chart, CSV / PowerPoint export.",
])

# ---- Slide 8: Requirement Specification ----
configure_body(body_shape(s[7]))
set_bullets(body_shape(s[7]), [
    "Functional — NL→SQL generation; deterministic introspection; classification + safety; RBAC; human review with dry-run / snapshot; multi-DB connect; AI analysis, dashboards and export.",
    "Non-functional — security: block DDL / injection, encrypt credentials (Fernet).",
    "Reliability — introspection always works; automatic cloud↔local failover.",
    "Performance — millisecond introspection, sub-second cloud generation.",
    "Usability, privacy (local mode) and extensibility (new engine = new adapter).",
])

# ---- Slide 9: Design Approach ----
configure_body(body_shape(s[8]))
set_bullets(body_shape(s[8]), [
    "Layered — presentation → application engines → LLM provider (failover) → data-access (Adapter) → persistence.",
    "Adapter pattern — one DatabaseAdapter interface (connect, get_schema, list_tables, execute) with eight concrete adapters; app code never imports a driver.",
    "A two-layer guardrail (classification + safety) plus an explicit human gate isolate the LLM from the database.",
    "Conversation context (last five turns) enables follow-up and query refinement.",
])

# ---- Slide 10: Architecture diagram (image only, clean 5-layer view) ----
remove(body_shape(s[9]))
add_image(s[9], f"{DIAG}/arch_ppt.png", 0.5, 1.3, 9.0, 3.7)

# ---- Slide 11: Innovations ----
configure_body(body_shape(s[10]))
set_bullets(body_shape(s[10]), [
    "Hybrid execution — deterministic introspection + LLM generation: correctness where cheap, intelligence where needed.",
    "One adapter abstraction extends NL querying across eight relational and NoSQL engines (JSON contract for Mongo / Redis).",
    "Cloud-to-local failover — resilient and privacy-flexible.",
    "Defence-in-depth — the LLM is never trusted to be safe; a separate validator and human gate guard execution.",
])

# ---- Slide 12: Implementation ----
configure_body(body_shape(s[11]))
set_bullets(body_shape(s[11]), [
    "Python 3.11 + Flask backend; two front-ends (Jinja + React 19 SPA) over one JSON API.",
    "LLM layer — Groq SDK (Llama 3.3 70B) and local Ollama (Mistral) with per-dialect prompt templates.",
    "Validator, Fernet-encrypted connection manager, and a per-engine snapshot / rollback engine.",
    "AI analysis via Groq JSON mode; dashboards persisted as JSON; PowerPoint export via python-pptx.",
])

# ---- Slide 13: Evaluation Metrics ----
configure_body(body_shape(s[12]))
set_bullets(body_shape(s[12]), [
    "Functional correctness — each feature works on the sample databases.",
    "Generation correctness — generated SQL executes and returns the intended rows.",
    "Safety — destructive and unauthorized statements are blocked.",
    "Latency — time to answer, by execution path and provider.",
    "Token efficiency — prompt and completion tokens per call.",
])

# ---- Slide 14: Testing & Validation ----
configure_body(body_shape(s[13]))
set_bullets(body_shape(s[13]), [
    "Unit — hardcoded commands (show tables = 12, describe, foreign keys) and the validator: all pass in < 20 ms.",
    "Integration — generation → validation → authorization → execution; write-review → snapshot → execute; failover: all pass.",
    "System — end-to-end NL query, DB switching, AI analysis, dashboards, destructive-command blocking: all pass on both front-ends.",
])

# ---- Slide 15: Performance Analysis (left bullets + chart right) ----
configure_body(body_shape(s[14]), top=1.45, height=3.4, left=0.4, width=4.5)
set_bullets(body_shape(s[14]), [
    "Groq cloud NL→SQL — mean 0.79 s, median 0.73 s (46 logged calls).",
    "Local Mistral — ~10.4 s (fully private, on-device).",
    "Hardcoded introspection — ~8 ms end-to-end (no LLM), ~100× faster.",
    "Generated SQL validated against expected results on the Chinook DB.",
])
add_image(s[14], f"{DIAG}/perf_hardcoded_vs_llm.png", 5.0, 1.5, 4.7, 3.3)

# ---- Slide 16: Hardware / System Requirements ----
configure_body(body_shape(s[15]))
set_bullets(body_shape(s[15]), [
    "Software — Python 3.11+ & Flask; Groq SDK / Ollama; DB drivers (psycopg2, PyMySQL, pymongo, redis, ...); cryptography; python-pptx; React / Vite / Tailwind.",
    "Hardware — any modern browser; 4 GB RAM (cloud) or 16 GB + optional GPU (local); 2–10 GB storage.",
    "Why — Flask is lightweight; Groq’s LPU gives low latency; Ollama enables privacy; per-engine drivers give multi-database reach.",
])

# ---- Slide 17: Conclusion ----
configure_body(body_shape(s[16]))
set_bullets(body_shape(s[16]), [
    "Meridian Data turns LLM text-to-SQL into a dependable tool — schema-aware prompting, deterministic shortcuts, layered validation, RBAC, human review and provider failover.",
    "Achieves sub-second cloud generation and millisecond introspection, and blocks every destructive operation.",
    "Serves as both an analyst productivity aid and a DBMS teaching instrument.",
])

# ---- Slide 18: Future Work ----
configure_body(body_shape(s[17]))
set_bullets(body_shape(s[17]), [
    "Replace the keyword safety filter with a full SQL parser (e.g. sqlglot).",
    "Unify the duplicated Jinja and JSON pipelines behind one service layer.",
    "Add retrieval-augmented schema selection for very large schemas.",
    "Self-correction — feed execution errors back to the model to repair queries.",
    "Stronger credential security; a quantitative accuracy study on Spider and BIRD.",
])

# ---- Slide 19: References (smaller font) ----
configure_body(body_shape(s[18]), top=1.4, height=3.6)
set_bullets(body_shape(s[18]), [
    "[1] E. F. Codd, “A relational model of data for large shared data banks,” Comm. ACM, 1970.",
    "[2] T. Yu et al., “Spider: cross-domain text-to-SQL dataset,” EMNLP, 2018.",
    "[3] T. Scholak et al., “PICARD: constrained auto-regressive decoding,” EMNLP, 2021.",
    "[4] M. Pourreza, D. Rafiei, “DIN-SQL: decomposed in-context learning,” NeurIPS, 2023.",
    "[5] J. Li et al., “BIRD: a big bench for database-grounded text-to-SQL,” NeurIPS, 2023.",
    "[6] R. S. Sandhu et al., “Role-based access control models,” IEEE Computer, 1996.",
], size=14)

# ---- Uniform, fit-guaranteed body font across content slides ----
for i in [1,2,3,4,5,6,7,8,10,11,12,13,14,15,16,17,18]:
    bs = body_shape(s[i])
    if bs is None: continue
    sz = 14 if i == 18 else 16
    for p in bs.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(sz)

# ===================== NEW IMAGE / TABLE SLIDES =====================
# Literature table
lit = duplicate(prs, 6); set_banner(lit, "Literature Survey — Summary")
remove(body_shape(lit))
rows_data = [
    ("Work", "Technique", "Key Limitation"),
    ("Seq2SQL (2017)", "RL over WikiSQL", "Single-table queries only"),
    ("SQLNet (2017)", "Sketch + column attention", "Tied to WikiSQL grammar"),
    ("Spider (2018)", "Cross-domain benchmark", "Schema-only; clean DBs"),
    ("PICARD (2021)", "Constrained decoding", "Syntactic, not semantic, validity"),
    ("DIN-SQL (2023)", "In-context decomposition", "High token cost / latency"),
    ("BIRD (2023)", "Dirty-DB benchmark", "GPT-4 only ~55% accuracy"),
]
tbl = lit.shapes.add_table(len(rows_data), 3, Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.4)).table
tbl.columns[0].width = Inches(2.2); tbl.columns[1].width = Inches(3.3); tbl.columns[2].width = Inches(3.3)
for ri, row in enumerate(rows_data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci); cell.text = val
        r = cell.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(14 if ri else 15); r.font.name = "Arial"; r.font.bold = (ri == 0)

# Methodology diagram
md = duplicate(prs, 6); set_banner(md, "Proposed Methodology — Pipeline")
remove(body_shape(md)); add_image(md, f"{DIAG}/methodology_v.png", 0.4, 1.7, 9.2, 2.9)

# ER diagram
er = duplicate(prs, 6); set_banner(er, "Database Design — ER Diagram (Sample Schema)")
remove(body_shape(er)); add_image(er, f"{DIAG}/er_chinook.png", 0.6, 1.35, 8.8, 3.7)

# Implementation screenshots
ss = duplicate(prs, 6); set_banner(ss, "Implementation — Working Application")
remove(body_shape(ss))
add_image(ss, f"{SCR}/styled_home.png", 0.4, 1.5, 4.6, 3.4)
add_image(ss, f"{SCR}/overview.png", 5.1, 1.5, 4.5, 3.4)

# Evaluation results (headline metrics + method/result)
ev = duplicate(prs, 6); set_banner(ev, "Evaluation Results")
add_image(ev, f"{DIAG}/eval_headline.png", 0.5, 1.15, 9.0, 2.25)
evb = body_shape(ev); configure_body(evb, top=3.55, height=1.5)
set_bullets(evb, [
    "Method — labelled NL→SQL queries on the Chinook DB (generated SQL executed and compared to gold result sets), a 12-prompt destructive / injection battery, and latency over 113 logged LLM calls.",
    "Result — correct SQL across aggregate, projection, top-N, group-by and join queries; 100% of unsafe statements contained (10 blocked, 2 routed to human review) with no false blocks on valid reads.",
])
for p in evb.text_frame.paragraphs:
    for r in p.runs: r.font.size = Pt(13)

# ===================== REORDER =====================
desired = [s[0],s[1],s[2],s[3],s[4],s[5], lit, s[6], md, s[7],s[8],s[9], er,
           s[10],s[11], ss, s[12],s[13],s[14], ev, s[15],s[16],s[17],s[18],s[19]]
sldIdLst = prs.slides._sldIdLst
current = list(sldIdLst)
mp = {id(prs.slides[i]._element): current[i] for i in range(len(current))}
for el in current: sldIdLst.remove(el)
for sl in desired: sldIdLst.append(mp[id(sl._element)])

prs.save(OUT)
print("WROTE", OUT, "slides:", len(list(prs.slides)))
